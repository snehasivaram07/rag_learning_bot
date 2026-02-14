
# import streamlit as st
# from transformers import pipeline
# from backend.loader import load_pdf_text
# from backend.chunker import chunk_text
# from backend.embeddings import embed_chunks, model
# from backend.vector_store import create_faiss_index
# from backend.rag_pipeline import answer_question

# from pptx import Presentation
# from docx import Document
# from pypdf import PdfReader
# import io


# # ---------------- FILE TEXT EXTRACTION ----------------

# def extract_text_from_file(uploaded_file):
#     file_type = uploaded_file.name.split(".")[-1].lower()

#     if file_type == "pdf":
#         reader = PdfReader(uploaded_file)
#         pages = []
#         for i, page in enumerate(reader.pages):
#             pages.append({
#                 "page": i + 1,
#                 "text": page.extract_text() or ""
#             })
#         return pages

#     elif file_type in ["ppt", "pptx"]:
#         prs = Presentation(uploaded_file)
#         pages = []
#         for i, slide in enumerate(prs.slides):
#             text = ""
#             for shape in slide.shapes:
#                 if hasattr(shape, "text"):
#                     text += shape.text + " "
#             pages.append({
#                 "page": i + 1,
#                 "text": text
#             })
#         return pages

#     elif file_type == "docx":
#         doc = Document(uploaded_file)
#         text = "\n".join([para.text for para in doc.paragraphs])
#         return [{"page": 1, "text": text}]

#     elif file_type == "txt":
#         stringio = io.StringIO(uploaded_file.getvalue().decode("utf-8"))
#         text = stringio.read()
#         return [{"page": 1, "text": text}]

#     else:
#         return None



# # ---------------- SUMMARIZER ----------------

# @st.cache_resource
# def load_summarizer():
#     return pipeline("summarization", model="facebook/bart-large-cnn")


# def generate_short_notes(text):
#     summarizer = load_summarizer()
#     text = text[:3000]
#     summary = summarizer(
#         text,
#         max_length=250,
#         min_length=120,
#         do_sample=False
#     )
#     return summary[0]["summary_text"]


# # ---------------- PAGE CONFIG ----------------

# st.set_page_config(layout="wide")
# st.title("📘 Study Chatbot")

# # ---------------- SESSION STATE ----------------

# if "messages" not in st.session_state:
#     st.session_state.messages = []

# if "planner" not in st.session_state:
#     st.session_state.planner = []

# # ---------------- SIDEBAR ----------------

# st.sidebar.header("📂 Upload Document")

# uploaded_file = st.sidebar.file_uploader(
#     "Upload file",
#     type=["pdf", "ppt", "pptx", "docx", "txt"]
# )

# if uploaded_file:
#     with st.spinner("Processing document..."):

#         pages = extract_text_from_file(uploaded_file)

#         if pages is None:
#             st.sidebar.error("Unsupported file type.")
#         else:
#             chunks = chunk_text(pages)
#             embeddings = embed_chunks(chunks)
#             index = create_faiss_index(embeddings)

#             st.session_state["document_text"] = pages
#             st.session_state["chunks"] = chunks
#             st.session_state["index"] = index

#             st.sidebar.success("Document processed successfully")

# # ---------------- DOCUMENT TOOLS ----------------

# st.sidebar.markdown("### 🛠 Document Tools")

# if st.sidebar.button("📝 Generate Short Notes"):
#     if "document_text" not in st.session_state:
#         st.sidebar.warning("Upload a document first.")
#     else:
#         with st.spinner("Generating short notes..."):
#             text = " ".join(page["text"] for page in st.session_state["document_text"])
#             notes = generate_short_notes(text)
#             st.session_state["short_notes"] = notes

# # ---------------- STUDY PLANNER ----------------

# st.sidebar.header("📅 Study Planner")
# topic = st.sidebar.text_input("Topic")
# date = st.sidebar.date_input("Date")

# if st.sidebar.button("Add"):
#     if topic:
#         st.session_state.planner.append(f"{date} → {topic}")

# for p in st.session_state.planner:
#     st.sidebar.write("•", p)

# # ---------------- SHORT NOTES DISPLAY ----------------

# if "short_notes" in st.session_state:
#     st.markdown("## 📝 Short Notes from Document")
#     st.info(st.session_state["short_notes"])

# # ---------------- CHAT HISTORY ----------------

# for msg in st.session_state.messages:
#     with st.chat_message(msg["role"]):
#         st.markdown(msg["content"])

# # ---------------- CHAT INPUT ----------------

# query = st.chat_input("Ask your question...")

# if query:
#     st.session_state.messages.append({"role": "user", "content": query})

#     with st.chat_message("user"):
#         st.markdown(query)

#     if "index" not in st.session_state:
#         answer_text = "Please upload a document first."
#     else:
#         response = answer_question(
#             query,
#             st.session_state["index"],
#             st.session_state["chunks"],
#             model
#         )

#         answer_text = response["answer"]

#         # Strict document mode — no hallucination
#         if response["mode"] == "document":
#             answer_text += "\n\n📌 **Sources:**\n"
#             for s in response["source"]:
#                 answer_text += f"- Page {s['page']}\n"
#         else:
#             answer_text = "The answer was not found clearly in the uploaded document."

#     with st.chat_message("assistant"):
#         st.markdown(answer_text)

#     st.session_state.messages.append({
#         "role": "assistant",
#         "content": answer_text
#     })


# import streamlit as st
# import io
# import tempfile
# from datetime import date

# # -------- Voice --------
# from streamlit_mic_recorder import mic_recorder
# from faster_whisper import WhisperModel
# from gtts import gTTS

# # -------- File Processing --------
# from pptx import Presentation
# from docx import Document
# from pypdf import PdfReader
# from pdf2image import convert_from_bytes
# import pytesseract

# # -------- Your Backend --------
# from backend.chunker import chunk_text
# from backend.embeddings import embed_chunks, model
# from backend.vector_store import create_faiss_index
# from backend.rag_pipeline import answer_question

# # =====================================================
# # PAGE CONFIG
# # =====================================================

# st.set_page_config(layout="wide")
# st.title("📘 Intelligent Study Chatbot")

# # =====================================================
# # LOAD MODELS (CACHED)
# # =====================================================

# @st.cache_resource
# def load_whisper():
#     model = WhisperModel("base", device="cpu", compute_type="int8")
#     return model

# whisper_model = load_whisper()


# # =====================================================
# # TEXT TO SPEECH
# # =====================================================

# def speak_text(text):
#     tts = gTTS(text)
#     tts.save("response.mp3")
#     audio_file = open("response.mp3", "rb")
#     st.audio(audio_file.read(), format="audio/mp3")

# # =====================================================
# # FILE EXTRACTION WITH OCR
# # =====================================================

# def extract_text_from_file(uploaded_file):
#     file_type = uploaded_file.name.split(".")[-1].lower()

#     # ---------- PDF ----------
#     if file_type == "pdf":
#         reader = PdfReader(uploaded_file)
#         pages = []

#         for i, page in enumerate(reader.pages):
#             text = page.extract_text()

#             # OCR fallback if text empty
#             if not text or len(text.strip()) < 10:
#                 images = convert_from_bytes(uploaded_file.getvalue())
#                 text = pytesseract.image_to_string(images[i])

#             pages.append({
#                 "page": i + 1,
#                 "text": text or ""
#             })

#         return pages

#     # ---------- PPT ----------
#     elif file_type in ["ppt", "pptx"]:
#         prs = Presentation(uploaded_file)
#         pages = []

#         for i, slide in enumerate(prs.slides):
#             text = ""
#             for shape in slide.shapes:
#                 if hasattr(shape, "text") and shape.text.strip():
#                     text += shape.text + " "

#             pages.append({
#                 "page": i + 1,
#                 "text": text
#             })

#         return pages

#     # ---------- DOCX ----------
#     elif file_type == "docx":
#         doc = Document(uploaded_file)
#         text = "\n".join([p.text for p in doc.paragraphs])
#         return [{"page": 1, "text": text}]

#     # ---------- TXT ----------
#     elif file_type == "txt":
#         text = uploaded_file.getvalue().decode("utf-8")
#         return [{"page": 1, "text": text}]

#     else:
#         return None


# # =====================================================
# # SESSION STATE
# # =====================================================

# if "messages" not in st.session_state:
#     st.session_state.messages = []

# if "planner" not in st.session_state:
#     st.session_state.planner = []

# # =====================================================
# # SIDEBAR
# # =====================================================

# st.sidebar.header("📂 Upload Document")

# uploaded_file = st.sidebar.file_uploader(
#     "Upload file",
#     type=["pdf", "ppt", "pptx", "docx", "txt"]
# )

# if uploaded_file:
#     with st.spinner("Processing document..."):
#         pages = extract_text_from_file(uploaded_file)

#         if pages and sum(len(p["text"]) for p in pages) > 50:
#             chunks = chunk_text(pages)
#             embeddings = embed_chunks(chunks)
#             index = create_faiss_index(embeddings)

#             st.session_state["document_text"] = pages
#             st.session_state["chunks"] = chunks
#             st.session_state["index"] = index

#             st.sidebar.success("Document processed successfully")
#         else:
#             st.sidebar.error("Could not extract readable text from file.")

# # =====================================================
# # STUDY PLANNER
# # =====================================================

# st.sidebar.header("📅 Study Planner")
# topic = st.sidebar.text_input("Topic")
# study_date = st.sidebar.date_input("Date")

# if st.sidebar.button("Add"):
#     if topic:
#         st.session_state.planner.append(f"{study_date} → {topic}")

# for p in st.session_state.planner:
#     st.sidebar.write("•", p)

# # =====================================================
# # CHAT HISTORY
# # =====================================================

# for msg in st.session_state.messages:
#     with st.chat_message(msg["role"]):
#         st.markdown(msg["content"])

# # =====================================================
# # VOICE INPUT
# # =====================================================

# st.markdown("### 🎙 Speak your question")

# audio = mic_recorder(
#     start_prompt="Start Recording",
#     stop_prompt="Stop Recording",
#     key="recorder"
# )

# voice_query = None

# if audio:
#     with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as tmp:
#         tmp.write(audio["bytes"])
#         tmp_path = tmp.name

#     segments, _ = whisper_model.transcribe(tmp_path)
# voice_query = ""
# for segment in segments:
#     voice_query += segment.text


# # =====================================================
# # CHAT INPUT
# # =====================================================

# query = st.chat_input("Ask your question...")

# if voice_query:
#     query = voice_query

# if query:
#     st.session_state.messages.append({"role": "user", "content": query})

#     with st.chat_message("user"):
#         st.markdown(query)

#     if "index" not in st.session_state:
#         answer_text = "Please upload a document first."
#     else:
#         response = answer_question(
#             query,
#             st.session_state["index"],
#             st.session_state["chunks"],
#             model
#         )

#         if response["mode"] == "document":
#             answer_text = response["answer"]
#             answer_text += "\n\n📌 **Sources:**\n"
#             for s in response["source"]:
#                 answer_text += f"- Page {s['page']}\n"
#         else:
#             answer_text = "The answer was not found clearly in the uploaded document."

#     with st.chat_message("assistant"):
#         st.markdown(answer_text)

#     # 🔊 Speak response
#     speak_text(answer_text)

#     st.session_state.messages.append({
#         "role": "assistant",
#         "content": answer_text
#     })


import streamlit as st
import io
import tempfile
from datetime import date

# ---------------- Voice ----------------
from streamlit_mic_recorder import mic_recorder
from faster_whisper import WhisperModel
from gtts import gTTS

# ---------------- File Processing ----------------
from pptx import Presentation
from docx import Document
from pypdf import PdfReader
from pdf2image import convert_from_bytes
import pytesseract

# ---------------- Your Backend ----------------
from backend.chunker import chunk_text
from backend.embeddings import embed_chunks, model
from backend.vector_store import create_faiss_index
from backend.rag_pipeline import answer_question

from transformers import pipeline


# =====================================================
# PAGE CONFIG
# =====================================================




st.set_page_config(layout="wide")
st.title("📘 Intelligent Study Chatbot")

# =====================================================
# LOAD WHISPER MODEL (WINDOWS SAFE)
# =====================================================

@st.cache_resource
def load_whisper():
    return WhisperModel(
        "base",
        device="cpu",
        compute_type="int8"   # Safe + lightweight
    )

whisper_model = load_whisper()

# =====================================================
# LOAD SUMMARIZER MODEL
# =====================================================

@st.cache_resource
def load_summarizer():
    return pipeline(
        "summarization",
        model="facebook/bart-large-cnn"
    )

summarizer = load_summarizer()


# =====================================================
# TEXT TO SPEECH
# =====================================================

def speak_text(text):
    tts = gTTS(text)
    tts.save("response.mp3")
    audio_file = open("response.mp3", "rb")
    st.audio(audio_file.read(), format="audio/mp3")

# =====================================================
# GENERATE SHORT NOTES
# =====================================================

def generate_short_notes(pages):

    full_text = " ".join([p["text"] for p in pages])

    max_chunk = 1000
    text_chunks = [
        full_text[i:i+max_chunk]
        for i in range(0, len(full_text), max_chunk)
    ]

    summaries = []

    for chunk in text_chunks:
        if len(chunk.strip()) > 50:
            summary = summarizer(
                "Create clear academic short notes with bullet points:\n\n" + chunk,
                max_length=200,
                min_length=80,
                do_sample=False
            )
            summaries.append(summary[0]["summary_text"])

    final_notes = "\n\n".join(summaries)
    return final_notes


# =====================================================
# FILE EXTRACTION WITH OCR
# =====================================================

def extract_text_from_file(uploaded_file):
    file_type = uploaded_file.name.split(".")[-1].lower()

    # ---------- PDF ----------
    if file_type == "pdf":
        reader = PdfReader(uploaded_file)
        pages = []

        for i, page in enumerate(reader.pages):
            text = page.extract_text()

            # OCR fallback if text missing
            if not text or len(text.strip()) < 10:
                images = convert_from_bytes(uploaded_file.getvalue())
                text = pytesseract.image_to_string(images[i])

            pages.append({
                "page": i + 1,
                "text": text or ""
            })

        return pages

    # ---------- PPT ----------
    elif file_type in ["ppt", "pptx"]:
        prs = Presentation(uploaded_file)
        pages = []

        for i, slide in enumerate(prs.slides):
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + " "

            pages.append({
                "page": i + 1,
                "text": text
            })

        return pages

    # ---------- DOCX ----------
    elif file_type == "docx":
        doc = Document(uploaded_file)
        text = "\n".join([p.text for p in doc.paragraphs])
        return [{"page": 1, "text": text}]

    # ---------- TXT ----------
    elif file_type == "txt":
        text = uploaded_file.getvalue().decode("utf-8")
        return [{"page": 1, "text": text}]

    return None


# =====================================================
# SESSION STATE
# =====================================================

if "messages" not in st.session_state:
    st.session_state.messages = []

if "planner" not in st.session_state:
    st.session_state.planner = []

# =====================================================
# SIDEBAR - FILE UPLOAD
# =====================================================

st.sidebar.header("📂 Upload Document")

uploaded_file = st.sidebar.file_uploader(
    "Upload file",
    type=["pdf", "ppt", "pptx", "docx", "txt"]
)

if uploaded_file:
    with st.spinner("Processing document..."):

        pages = extract_text_from_file(uploaded_file)

        if pages and sum(len(p["text"]) for p in pages) > 50:
            chunks = chunk_text(pages)
            embeddings = embed_chunks(chunks)
            index = create_faiss_index(embeddings)

            st.session_state["document_text"] = pages
            st.session_state["chunks"] = chunks
            st.session_state["index"] = index

            st.sidebar.success("Document processed successfully")
        else:
            st.sidebar.error("Could not extract readable text from file.")

# =====================================================
# STUDY PLANNER
# =====================================================

st.sidebar.header("📅 Study Planner")
topic = st.sidebar.text_input("Topic")
study_date = st.sidebar.date_input("Date")

if st.sidebar.button("Add"):
    if topic:
        st.session_state.planner.append(f"{study_date} → {topic}")

for p in st.session_state.planner:
    st.sidebar.write("•", p)

# =====================================================
# SHORT NOTES SECTION
# =====================================================

st.sidebar.header("📝 Generate Short Notes")

if st.sidebar.button("Create Short Notes"):
    if "document_text" in st.session_state:
        with st.spinner("Generating short notes..."):
            notes = generate_short_notes(
                st.session_state["document_text"]
            )

            st.session_state["short_notes"] = notes

            st.success("Short notes generated successfully!")
    else:
        st.warning("Please upload a document first.")



# =====================================================
# DISPLAY CHAT HISTORY
# =====================================================

for msg in st.session_state.messages:
    with st.chat_message(msg["role"]):
        st.markdown(msg["content"])
# =====================================================
# DISPLAY SHORT NOTES
# =====================================================

if "short_notes" in st.session_state:
    st.subheader("📌 Short Notes")
    st.markdown(st.session_state["short_notes"])

# =====================================================
# VOICE INPUT
# =====================================================

st.markdown("### 🎙 Speak your question")

audio = mic_recorder(
    start_prompt="Start Recording",
    stop_prompt="Stop Recording",
    key="recorder"
)

voice_query = None

if audio:
    with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as tmp:
        tmp.write(audio["bytes"])
        tmp_path = tmp.name

    segments, _ = whisper_model.transcribe(tmp_path)

    voice_query = ""
    for segment in segments:
        voice_query += segment.text

    st.write("You said:", voice_query)

# =====================================================
# CHAT INPUT
# =====================================================

query = st.chat_input("Ask your question...")

if voice_query:
    query = voice_query

if query:
    st.session_state.messages.append({"role": "user", "content": query})

    with st.chat_message("user"):
        st.markdown(query)

    if "index" not in st.session_state:
        answer_text = "Please upload a document first."
    else:
        response = answer_question(
            query,
            st.session_state["index"],
            st.session_state["chunks"],
            model
        )

        if response["mode"] == "document":
            answer_text = response["answer"]
            answer_text += "\n\n📌 **Sources:**\n"
            for s in response["source"]:
                answer_text += f"- Page {s['page']}\n"
        else:
            answer_text = "The answer was not found clearly in the uploaded document."

    with st.chat_message("assistant"):
        st.markdown(answer_text)

    # 🔊 Voice response
    speak_text(answer_text)

    st.session_state.messages.append({
        "role": "assistant",
        "content": answer_text
    })
